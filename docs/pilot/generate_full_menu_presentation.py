from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_FILE = BASE_DIR / "Semantra_Full_Menu_Presentation_2026-05-27.pptx"

ACCENT = RGBColor(17, 64, 124)
TEXT = RGBColor(34, 34, 34)
MUTED = RGBColor(96, 96, 96)
BG = RGBColor(248, 248, 246)
PANEL = RGBColor(234, 240, 248)


@dataclass
class SlideSpec:
    title: str
    subtitle: str
    key_message: str
    evidence: list[str] = field(default_factory=list)
    speaker_note: str = ""
    image_paths: list[Path] = field(default_factory=list)


def existing_paths(*relative_paths: str) -> list[Path]:
    results: list[Path] = []
    for value in relative_paths:
        candidate = Path(value)
        if not candidate.is_absolute():
            candidate = BASE_DIR.parent / value
        if candidate.exists():
            results.append(candidate)
    return results


def build_specs() -> list[SlideSpec]:
    assets_root = BASE_DIR / "demo_assets"
    workspace_assets = assets_root / "workspace_recordings_20260527"
    manual_assets = assets_root / "manual_live_demo_20260527"
    supporting_assets = assets_root / "full_menu_supporting_assets_20260527"

    return [
        SlideSpec(
            title="Semantra",
            subtitle="Full menu presentation",
            key_message="Semantra povezuje analyst rad, quality evidence i governance odluke u jedinstven data-integration tok.",
            evidence=[
                "Radni tok: Workspace -> Catalog -> Benchmarks -> System -> Governance",
                "Pokazuje ingestion, review, decisions, output, reuse, observability i stewardship",
            ],
            speaker_note="Semantra nije samo mapper. To je kompletan radni i governance sistem za data integration, od pripreme i review-a do benchmark merenja, runtime kontrole i stewardship odluka.",
            image_paths=existing_paths(
                str(workspace_assets / "04_workspace_output_generation" / "screenshots" / "04_workspace_output_generation_01.png")
            ),
        ),
        SlideSpec(
            title="Navigation Map",
            subtitle="Top-level product structure",
            key_message="Pet glavnih zona razdvajaju analyst rad, reuse, quality telemetry, runtime kontrolu i steward odgovornosti.",
            evidence=[
                "Workspace: setup, review, decisions, output",
                "Catalog: discovery, detail, diff, reuse, handoff",
                "Benchmarks, System i Governance pokrivaju kvalitet, operacije i stewarding",
            ],
            speaker_note="Na vrhu aplikacije postoji pet glavnih zona: Workspace, Catalog, Benchmarks, System i Governance. Time publika odmah dobija mentalni model cele aplikacije.",
            image_paths=existing_paths(
                str(workspace_assets / "04_workspace_output_generation" / "screenshots" / "04_workspace_output_generation_01.png")
            ),
        ),
        SlideSpec(
            title="Workspace / Setup",
            subtitle="Upload and interpret source inputs",
            key_message="Setup pokriva standardni two-file mapping i canonical-only tok, uz razlikovanje row-data i schema-spec inputa.",
            evidence=[
                "Mapping mode: Standard ili Canonical",
                "Source/Target upload i profiling",
                "Source mode i Target mode za data vs. schema spec",
            ],
            speaker_note="Setup je mesto gde pocinje rad sa podacima. Ovde korisnik bira da li radi standardni two-file mapping ili canonical-only tok, i ovde sistem tumaci da li je fajl row data ili schema spec.",
            image_paths=existing_paths(
                str(workspace_assets / "01_standard_two_file_mapping" / "screenshots" / "01_standard_two_file_mapping_01.png")
            ),
        ),
        SlideSpec(
            title="Workspace / Review",
            subtitle="Trust-layer analyst review",
            key_message="Review prikazuje kandidate, coverage, knowledge i LLM signale, ali zadrzava bounded analyst-centered tok odluka.",
            evidence=[
                "Mapping Trust Layer i review prioritizacija",
                "Review Queue Plan za fokus rada",
                "LLM Decision Proposals kao bounded pomocni sloj",
            ],
            speaker_note="Review je analyst-centered trust layer. Tu se vide kandidati, coverage, knowledge i LLM signali, kao i review prioritizacija kroz queue plan.",
            image_paths=existing_paths(
                str(workspace_assets / "03_llm_decision_flow" / "screenshots" / "03_llm_decision_flow_01.png")
            ),
        ),
        SlideSpec(
            title="Workspace / Decisions",
            subtitle="Persist and resume analyst decisions",
            key_message="Decisions pretvara review state u trajne ili polutrajne odluke kroz overrides, imports, versions i draft sessions.",
            evidence=[
                "Manual overrides i import/export povrsine",
                "Mapping Set Versions i draft-session restore",
                "Continuity rada bez gubitka konteksta",
            ],
            speaker_note="Decisions je prelaz iz review-a u trajne ili polutrajne odluke. Tu se cuvaju manual overrides, import/export stanja, mapping set verzije, draft sessions i correction tokovi.",
            image_paths=existing_paths(
                str(manual_assets / "screenshots" / "02_workspace_resume_01.png"),
                str(workspace_assets / "03_llm_decision_flow" / "screenshots" / "03_llm_decision_flow_01.png"),
            ),
        ),
        SlideSpec(
            title="Workspace / Output",
            subtitle="From mapping state to dev artifacts",
            key_message="Output zatvara analyst tok kroz preview, Pandas, PySpark, dbt i LLM refinement artefakte.",
            evidence=[
                "Preview za standardni mapping",
                "Artifact format za Pandas, PySpark i dbt",
                "Refinement kao dodatni polish korak",
            ],
            speaker_note="Output pretvara aktivni mapping state u razvojne artefakte. Iz istog radnog stanja mogu da nastanu preview, Pandas, PySpark ili dbt izlazi, plus refinement kada je potreban.",
            image_paths=existing_paths(
                str(workspace_assets / "04_workspace_output_generation" / "screenshots" / "04_workspace_output_generation_01.png"),
                str(workspace_assets / "04_workspace_output_generation" / "screenshots" / "04_workspace_output_generation_02.png"),
            ),
        ),
        SlideSpec(
            title="Catalog / Search and Discovery",
            subtitle="Reusable integration knowledge",
            key_message="Catalog omogucava discovery reusable integracija po sistemima, domenu, statusu i canonical kontekstu.",
            evidence=[
                "Search and Filters za system/domain/owner/status",
                "Discovery Overview i Integration Results",
                "Reuse shortlist i system-pair pregled",
            ],
            speaker_note="Catalog je reuse i discovery biblioteka integracionog znanja. Pretraga radi po sistemima, domenu, statusu, owner-u i canonical signalima, pa korisnik ne mora da krece od nule.",
            image_paths=existing_paths(
                str(manual_assets / "screenshots" / "01_catalog_reuse_01.png")
            ),
        ),
        SlideSpec(
            title="Catalog / Detail, Diff and Handoff",
            subtitle="Move from stored assets into live work",
            key_message="Catalog nije pasivan: iz njega se ulazi u Workspace reuse, diff review i governance handoff tokove.",
            evidence=[
                "Approved reuse u Workspace",
                "Version diff i Workspace Review handoff",
                "Stewardship handoff u Governance",
            ],
            speaker_note="Kada otvorimo konkretan asset, Catalog postaje ulaz u sledeci korak rada. Iz njega mozemo da uradimo reuse u Workspace, diff handoff u Workspace Review ili governance handoff u Stewardship.",
            image_paths=existing_paths(
                str(manual_assets / "screenshots" / "01_catalog_reuse_01.png"),
                str(manual_assets / "screenshots" / "03_catalog_diff_handoff_01.png"),
                str(manual_assets / "screenshots" / "04_catalog_stewardship_handoff_01.png"),
            ),
        ),
        SlideSpec(
            title="Benchmarks / Datasets and Runs",
            subtitle="Persistent quality evidence",
            key_message="Benchmarking pocinje od cuvanja mapping stanja kao evaluation dataset-a i vracanja istorije run-ova.",
            evidence=[
                "Save current mapping as benchmark",
                "Load saved benchmark datasets",
                "Load benchmark run history",
            ],
            speaker_note="Benchmarks pocinju od cuvanja mapping stanja kao benchmark dataset-a i od ucitavanja prethodnih run-ova. To znaci da evaluacija nije jednokratna demonstracija, vec trajni kvalitetni signal.",
            image_paths=existing_paths(
                str(manual_assets / "screenshots" / "05_benchmarks_explanation_01.png")
            ),
        ),
        SlideSpec(
            title="Benchmarks / Profile Comparison",
            subtitle="Explainable scoring recommendations",
            key_message="Benchmarks ne daju samo metriku, vec i preporuceni scoring profil sa objasnjenjem, rizicima i sledecim koracima.",
            evidence=[
                "Scoring Profile Comparison",
                "Recommended default profile",
                "Benchmark Explanation sa findings i risks",
            ],
            speaker_note="Kada benchmark dataset postoji, mozemo da uporedimo vise scoring profila i da dobijemo preporuceni default. Posle toga explanation prevodi metrike u razumljiv poslovni razlog.",
            image_paths=existing_paths(
                str(manual_assets / "screenshots" / "05_benchmarks_explanation_01.png")
            ),
        ),
        SlideSpec(
            title="System / Admin",
            subtitle="Runtime operations and control",
            key_message="System Admin je kontrolna tabla za runtime config, scoring profil, corrections i benchmark-run administraciju.",
            evidence=[
                "Load runtime config, saved corrections i benchmark runs",
                "Scoring Runtime sa aktivnim profilom",
                "Apply scoring profile za nove mapping prolaze",
            ],
            speaker_note="System Admin je runtime administracija. Tu se ucitava runtime config, vide saved corrections i benchmark runs, i menja aktivni scoring profil za nove mapping prolaze.",
            image_paths=existing_paths(
                str(supporting_assets / "11_system_admin_01.png")
            ),
        ),
        SlideSpec(
            title="System / Debug",
            subtitle="Structured observability",
            key_message="System Debug objedinjuje decision logs, knowledge runtime status, audit log i coverage debug pogled nad aktivnim mapping stanjem.",
            evidence=[
                "Decision logs i knowledge runtime status",
                "Knowledge audit log",
                "Coverage and match insights za troubleshooting",
            ],
            speaker_note="System Debug je observability povrsina. Tu se vide decision logs, aktivni knowledge runtime, audit log i canonical coverage/debug insighte nad trenutnim mapping stanjem.",
            image_paths=existing_paths(
                str(supporting_assets / "12_system_debug_01.png")
            ),
        ),
        SlideSpec(
            title="Governance Overview",
            subtitle="Steward console structure",
            key_message="Governance razdvaja cetiri odgovornosti: Canonical, Knowledge, Overlays & Runtime i Stewardship.",
            evidence=[
                "Canonical za stabilni glossary",
                "Knowledge za radni registry pojmova",
                "Overlays & Runtime i Stewardship za operativno upravljanje promenama",
            ],
            speaker_note="Governance je steward konzola sa cetiri sekcije: Canonical, Knowledge, Overlays & Runtime i Stewardship. Ovaj slajd sluzi da publika shvati da governance nije jedan ekran.",
            image_paths=existing_paths(
                str(supporting_assets / "16_governance_overlays_runtime_01.png")
            ),
        ),
        SlideSpec(
            title="Governance / Canonical",
            subtitle="Stable glossary stewardship",
            key_message="Canonical sekcija upravlja stabilnim konceptima, aliasima i promotion-ready signalima koji imaju dugorocan efekat.",
            evidence=[
                "Canonical Glossary",
                "Concept details, aliases i context coverage",
                "Promotion-oriented governance tokovi",
            ],
            speaker_note="Canonical je stabilni glossary sloj. Tu se upravlja canonical konceptima, aliasima, coverage kontekstom i promotion-ready signalima.",
            image_paths=existing_paths(
                str(supporting_assets / "14_governance_canonical_01.png")
            ),
        ),
        SlideSpec(
            title="Governance / Knowledge",
            subtitle="Operational knowledge registry",
            key_message="Knowledge sekcija povezuje realne integracione izraze sa canonical slojem i omogucava postepenu steward normalizaciju.",
            evidence=[
                "Knowledge Registry i Knowledge Concept Registry",
                "Linked canonical paths",
                "Promotion readiness i review detalji",
            ],
            speaker_note="Knowledge cuva radni registry pojmova blizih stvarnim integracionim izrazima. On povezuje operativno znanje sa canonical slojem.",
            image_paths=existing_paths(
                str(supporting_assets / "15_governance_knowledge_01.png")
            ),
        ),
        SlideSpec(
            title="Governance / Overlays & Runtime",
            subtitle="Reversible runtime knowledge changes",
            key_message="Overlays & Runtime obezbedjuje brz, reverzibilan i auditabilan sloj za runtime promene bez direktnog menjanja stabilnog glossarya.",
            evidence=[
                "Overlay Summary i active overlay status",
                "Overlay lifecycle: refresh, reload, audit",
                "Runtime-aware knowledge management",
            ],
            speaker_note="Overlays & Runtime sluzi za kontrolisane i reverzibilne promene znanja koje uticu na runtime bez direktnog menjanja stabilnog glossarya.",
            image_paths=existing_paths(
                str(supporting_assets / "16_governance_overlays_runtime_01.png")
            ),
        ),
        SlideSpec(
            title="Governance / Stewardship",
            subtitle="Queue-based follow-up and decisioning",
            key_message="Stewardship pretvara operativne signale i gapove u auditabilne approve/reject/ignore odluke.",
            evidence=[
                "Queue status i selected item detail",
                "Review note i proposal-state signal",
                "Approve/reject/ignore tokovi za governance follow-up",
            ],
            speaker_note="Stewardship je radna lista governance follow-up stavki. Tu steward vidi gapove, review note, queue statuse i odluke za promotion, reject ili ignore.",
            image_paths=existing_paths(
                str(manual_assets / "screenshots" / "04_catalog_stewardship_handoff_01.png")
            ),
        ),
        SlideSpec(
            title="Closing",
            subtitle="One integrated operating model",
            key_message="Semantra povezuje ingestion, review, decisions, output, reuse, benchmark evidence, runtime kontrolu i governance upravljanje u jedan tok.",
            evidence=[
                "Catalog vraca korisnika u Workspace",
                "Benchmarks objasnjavaju kvalitet odluka",
                "Governance pretvara operativne signale u trajno upravljano znanje",
            ],
            speaker_note="Zakljucak prezentacije je da Semantra pokriva ceo zivotni ciklus: ingestion, review, decisions, output, reuse, benchmark evidence, runtime kontrolu i governance upravljanje znanjem.",
            image_paths=existing_paths(
                str(supporting_assets / "16_governance_overlays_runtime_01.png")
            ),
        ),
    ]


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(5.2), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(5.6), Inches(0.5))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.name = "Aptos"
    run2.font.size = Pt(12)
    run2.font.color.rgb = MUTED


def add_key_message(slide, text: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.75), Inches(5.7), Inches(0.95))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = PANEL
    shape.line.color.rgb = PANEL
    tf = shape.text_frame
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = ACCENT


def add_evidence_panel(slide, evidence: Iterable[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.95), Inches(5.6), Inches(3.55))
    tf = box.text_frame
    tf.word_wrap = True

    header = tf.paragraphs[0]
    header_run = header.add_run()
    header_run.text = "Why This Slide Matters"
    header_run.font.name = "Aptos"
    header_run.font.size = Pt(16)
    header_run.font.bold = True
    header_run.font.color.rgb = TEXT

    for item in evidence:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Aptos"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT


def add_image_collage(slide, image_paths: list[Path]) -> None:
    if not image_paths:
        return
    left = Inches(6.45)
    top = Inches(0.75)
    width = Inches(6.2)
    if len(image_paths) == 1:
        slide.shapes.add_picture(str(image_paths[0]), left, top, width=width)
        return
    if len(image_paths) == 2:
        slide.shapes.add_picture(str(image_paths[0]), left, top, width=width)
        slide.shapes.add_picture(str(image_paths[1]), left, Inches(3.95), width=width)
        return
    # 3+ images: one large at top, two small below
    slide.shapes.add_picture(str(image_paths[0]), left, top, width=width)
    slide.shapes.add_picture(str(image_paths[1]), left, Inches(4.15), width=Inches(3.0))
    slide.shapes.add_picture(str(image_paths[2]), Inches(9.65), Inches(4.15), width=Inches(3.0))


def add_footer(slide, note: str, index: int, total: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.45))
    tf = footer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Speaker note: {note}"
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.italic = True
    run.font.color.rgb = MUTED

    page_box = slide.shapes.add_textbox(Inches(12.4), Inches(6.82), Inches(0.7), Inches(0.3))
    ptf = page_box.text_frame
    pp = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    prun = pp.add_run()
    prun.text = f"{index}/{total}"
    prun.font.name = "Aptos"
    prun.font.size = Pt(10)
    prun.font.color.rgb = MUTED


def build_presentation(specs: list[SlideSpec]) -> Presentation:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    blank_layout = presentation.slide_layouts[6]
    total = len(specs)
    for index, spec in enumerate(specs, start=1):
        slide = presentation.slides.add_slide(blank_layout)
        add_background(slide)
        add_title(slide, spec.title, spec.subtitle)
        add_key_message(slide, spec.key_message)
        add_evidence_panel(slide, spec.evidence)
        add_image_collage(slide, spec.image_paths)
        add_footer(slide, spec.speaker_note, index, total)
    return presentation


def main() -> None:
    specs = build_specs()
    presentation = build_presentation(specs)
    presentation.save(str(OUTPUT_FILE))
    print(OUTPUT_FILE)


if __name__ == "__main__":
    main()